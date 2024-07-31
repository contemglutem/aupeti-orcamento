import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
import textwrap


class Table:

    def __init__(self):
        self.output_path = "/home/simon/PycharmProjects/orcamentoAupetitoso/data/3.pptx"
        self.output_png = "/home/simon/PycharmProjects/orcamentoAupetitoso/data/3.png"
        self.lista = ['Peito de frango (95 g)', 'Coração de frango (34 g)', 'Brócolis (60 g)', '-', '-']
        self.lista_maior = ['Peito de frango (95 g)',
                            'Coração de frango (34 g)','Peito de frango (95 g)',
                            'Coração de frango (34 g)','Peito de frango (95 g)',
                            'Coração de frango (34 g)','Peito de frango (95 g)',
                            'Coração de frango (34 g)','Peito de frango (95 g)',]
        
        self.dataframe = pd.DataFrame(data={
            'Receita 1': self.lista_maior,
            'Receita 2': self.lista_maior,
            'Receita 3': self.lista_maior,
        })
        self.presentation_path = self.select_presentation()

    def select_presentation(self):
        if self.dataframe.shape[1] == 1:
            return "/home/simon/PycharmProjects/orcamentoAupetitoso/lib/3_1.pptx"

        elif self.dataframe.shape[1] == 2:
            return "/home/simon/PycharmProjects/orcamentoAupetitoso/lib/3_2.pptx"

        elif self.dataframe.shape[1] == 3:
            return "/home/simon/PycharmProjects/orcamentoAupetitoso/lib/3_3.pptx"

    def convert_to_png(self):
        subprocess.call(['libreoffice', '--headless', '--convert-to', 'png', self.output_path, '--outdir',
                         '/home/simon/PycharmProjects/orcamentoAupetitoso/data'])

    def main(self):
        df = self.dataframe
        # Open the presentation
        prs = Presentation(self.presentation_path)
        slide = prs.slides[0]

        # Add table to slide
        x, y, cx, cy = (Inches(1.5/Table().dataframe.shape[1]),
                        Inches(3) - Table().dataframe.shape[0] * Inches(0.5) / 4,
                        Inches(Table().dataframe.shape[1]*2),
                        Inches(0.5))
        print(x, y, cx, cy)
        shape = slide.shapes.add_table(df.shape[0] + 1, df.shape[1], x, y, cx, cy)
        table = shape.table

        # Set table headers
        for col_index, col_name in enumerate(df.columns):
            cell = table.cell(0, col_index)
            cell.text = col_name
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Fill the table with DataFrame values
        for row_index, row in df.iterrows():
            for col_index, value in enumerate(row):
                cell = table.cell(row_index + 1, col_index)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Set the first row cells to green and white text
        for row_index in range(0, 1):
            for col_index in range(len(table.columns)):
                cell = table.cell(row_index, col_index)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 102, 54)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.name = "Glacial Indifference"
                # Remove shadow
                cell.text_frame.paragraphs[0].font.shadow = None

        # Set all other rows to red
        for row_index in range(1, len(table.rows)):
            for col_index in range(len(table.columns)):
                cell = table.cell(row_index, col_index)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 224, 213)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        run.font.name = "Glacial Indifference"
                # Remove shadow
                cell.text_frame.paragraphs[0].font.shadow = None

        # Save the presentation
        prs.save(self.output_path)

        # Convert to PNG using libreoffice
        self.convert_to_png()


if __name__ == "__main__":
    Table().main()
